#!/usr/bin/env python3
"""
Create a PowerPoint brand deck for Lonesome No More
Optimized for Gamma brand import
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor as RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Brand colors (hex to RGB)
COLORS = {
    'primary': RgbColor(0x5F, 0x7A, 0x61),      # Sage green
    'secondary': RgbColor(0xD4, 0x73, 0x5E),    # Terra cotta
    'accent': RgbColor(0xE8, 0xAA, 0x6B),       # Honey amber
    'text_dark': RgbColor(0x2C, 0x2C, 0x2C),    # Deep charcoal
    'text_medium': RgbColor(0x5A, 0x5A, 0x5A),  # Warm gray
    'bg_cream': RgbColor(0xFB, 0xF8, 0xF3),     # Warm cream
    'bg_light': RgbColor(0xFF, 0xFE, 0xFB),     # Soft white
    'white': RgbColor(0xFF, 0xFF, 0xFF),
}

def add_title_slide(prs, title, subtitle=""):
    """Add a title slide with brand styling"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Background color
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['bg_cream']
    background.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Merriweather"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.name = "Inter"
        p.font.size = Pt(24)
        p.font.color.rgb = COLORS['text_medium']
        p.alignment = PP_ALIGN.CENTER

    return slide


def add_color_swatch(slide, x, y, width, height, color, label, hex_code):
    """Add a color swatch with label"""
    # Color rectangle
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

    # Label below
    label_box = slide.shapes.add_textbox(x, y + height + Inches(0.1), width, Inches(0.4))
    tf = label_box.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.name = "Inter"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS['text_dark']
    p.alignment = PP_ALIGN.CENTER

    # Hex code
    hex_box = slide.shapes.add_textbox(x, y + height + Inches(0.45), width, Inches(0.3))
    tf = hex_box.text_frame
    p = tf.paragraphs[0]
    p.text = hex_code
    p.font.name = "Inter"
    p.font.size = Pt(12)
    p.font.color.rgb = COLORS['text_medium']
    p.alignment = PP_ALIGN.CENTER


def add_color_palette_slide(prs):
    """Add color palette slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['bg_light']
    background.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Color Palette"
    p.font.name = "Merriweather"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    p.alignment = PP_ALIGN.CENTER

    # Primary colors row
    colors_row1 = [
        (COLORS['primary'], "Primary\nSage Green", "#5F7A61"),
        (COLORS['secondary'], "Secondary\nTerra Cotta", "#D4735E"),
        (COLORS['accent'], "Accent\nHoney Amber", "#E8AA6B"),
    ]

    swatch_width = Inches(2.5)
    swatch_height = Inches(1.5)
    start_x = Inches(1)
    start_y = Inches(1.3)
    gap = Inches(0.5)

    for i, (color, label, hex_code) in enumerate(colors_row1):
        x = start_x + i * (swatch_width + gap)
        add_color_swatch(slide, x, start_y, swatch_width, swatch_height, color, label, hex_code)

    # Text and background colors row
    colors_row2 = [
        (COLORS['text_dark'], "Text Dark", "#2C2C2C"),
        (COLORS['text_medium'], "Text Medium", "#5A5A5A"),
        (COLORS['bg_cream'], "Background\nCream", "#FBF8F3"),
    ]

    start_y2 = Inches(4)
    for i, (color, label, hex_code) in enumerate(colors_row2):
        x = start_x + i * (swatch_width + gap)
        add_color_swatch(slide, x, start_y2, swatch_width, swatch_height, color, label, hex_code)

    return slide


def add_typography_slide(prs):
    """Add typography showcase slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['bg_cream']
    background.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Typography"
    p.font.name = "Merriweather"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    p.alignment = PP_ALIGN.CENTER

    # Display Font Section
    display_label = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.4))
    tf = display_label.text_frame
    p = tf.paragraphs[0]
    p.text = "Display Font: Merriweather"
    p.font.name = "Inter"
    p.font.size = Pt(14)
    p.font.color.rgb = COLORS['text_medium']

    # Merriweather samples
    y_pos = Inches(1.6)
    samples = [
        ("Merriweather", Pt(44), True, "Headline Text"),
        ("Merriweather", Pt(32), True, "Section Title"),
        ("Merriweather", Pt(24), False, "Subsection Heading"),
    ]

    for font_name, size, bold, text in samples:
        box = slide.shapes.add_textbox(Inches(0.5), y_pos, Inches(9), Inches(0.6))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = font_name
        p.font.size = size
        p.font.bold = bold
        p.font.color.rgb = COLORS['text_dark']
        y_pos += Inches(0.7)

    # Body Font Section
    body_label = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(0.4))
    tf = body_label.text_frame
    p = tf.paragraphs[0]
    p.text = "Body Font: Inter"
    p.font.name = "Inter"
    p.font.size = Pt(14)
    p.font.color.rgb = COLORS['text_medium']

    # Inter samples
    body_samples = [
        (Pt(19), False, "Body text for optimal readability (19px base)"),
        (Pt(16), False, "Secondary text and descriptions"),
        (Pt(14), True, "Button labels and UI elements"),
    ]

    y_pos = Inches(4.4)
    for size, bold, text in body_samples:
        box = slide.shapes.add_textbox(Inches(0.5), y_pos, Inches(9), Inches(0.5))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = "Inter"
        p.font.size = size
        p.font.bold = bold
        p.font.color.rgb = COLORS['text_dark']
        y_pos += Inches(0.5)

    return slide


def add_button_styles_slide(prs):
    """Add button component showcase"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['bg_light']
    background.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Button Styles"
    p.font.name = "Merriweather"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    p.alignment = PP_ALIGN.CENTER

    # Primary Button (Sage)
    btn1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.5), Inches(3), Inches(0.8))
    btn1.fill.solid()
    btn1.fill.fore_color.rgb = COLORS['primary']
    btn1.line.fill.background()
    tf = btn1.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = "Primary Button"
    p.font.name = "Inter"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    tf.anchor = MSO_ANCHOR.MIDDLE

    # CTA Button (Terra Cotta)
    btn2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.5), Inches(1.5), Inches(3), Inches(0.8))
    btn2.fill.solid()
    btn2.fill.fore_color.rgb = COLORS['secondary']
    btn2.line.fill.background()
    tf = btn2.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = "Get Started"
    p.font.name = "Inter"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    tf.anchor = MSO_ANCHOR.MIDDLE

    # Accent Button
    btn3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(2.8), Inches(3), Inches(0.8))
    btn3.fill.solid()
    btn3.fill.fore_color.rgb = COLORS['accent']
    btn3.line.fill.background()
    tf = btn3.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = "Learn More"
    p.font.name = "Inter"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    tf.anchor = MSO_ANCHOR.MIDDLE

    # Outlined Button
    btn4 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.5), Inches(2.8), Inches(3), Inches(0.8))
    btn4.fill.solid()
    btn4.fill.fore_color.rgb = COLORS['white']
    btn4.line.color.rgb = COLORS['primary']
    btn4.line.width = Pt(2)
    tf = btn4.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = "Secondary Button"
    p.font.name = "Inter"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    p.alignment = PP_ALIGN.CENTER
    tf.anchor = MSO_ANCHOR.MIDDLE

    # Card example
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(4.2), Inches(7), Inches(2))
    card.fill.solid()
    card.fill.fore_color.rgb = COLORS['white']
    card.line.color.rgb = RgbColor(0xD1, 0xD5, 0xDB)
    card.line.width = Pt(1)

    card_title = slide.shapes.add_textbox(Inches(1.8), Inches(4.4), Inches(6.5), Inches(0.5))
    tf = card_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Card Component"
    p.font.name = "Merriweather"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']

    card_text = slide.shapes.add_textbox(Inches(1.8), Inches(4.9), Inches(6.5), Inches(1))
    tf = card_text.text_frame
    p = tf.paragraphs[0]
    p.text = "Cards use white backgrounds with subtle borders and rounded corners for a clean, accessible design."
    p.font.name = "Inter"
    p.font.size = Pt(14)
    p.font.color.rgb = COLORS['text_medium']

    return slide


def add_logo_slide(prs, logo_path):
    """Add logo showcase slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['bg_cream']
    background.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Logo"
    p.font.name = "Merriweather"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    p.alignment = PP_ALIGN.CENTER

    # Add logo if it exists
    if os.path.exists(logo_path):
        slide.shapes.add_picture(logo_path, Inches(3.5), Inches(1.5), Inches(3), Inches(3))

    # Brand name
    brand_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(9), Inches(0.8))
    tf = brand_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Lonesome No More"
    p.font.name = "Merriweather"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    p.alignment = PP_ALIGN.CENTER

    # Tagline
    tag_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.6), Inches(9), Inches(0.5))
    tf = tag_box.text_frame
    p = tf.paragraphs[0]
    p.text = "AI Companionship for Seniors"
    p.font.name = "Inter"
    p.font.size = Pt(20)
    p.font.color.rgb = COLORS['text_medium']
    p.alignment = PP_ALIGN.CENTER

    return slide


def add_sample_content_slide(prs):
    """Add sample content layout slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['bg_light']
    background.line.fill.background()

    # Header bar
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = COLORS['primary']
    header.line.fill.background()

    # Header text
    header_text = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(4), Inches(0.5))
    tf = header_text.text_frame
    p = tf.paragraphs[0]
    p.text = "Lonesome No More"
    p.font.name = "Merriweather"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']

    # Main title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Companionship for Your\nLoved Ones"
    p.font.name = "Merriweather"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['text_dark']
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(0.6))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Peace of mind for you through simple phone conversations"
    p.font.name = "Inter"
    p.font.size = Pt(18)
    p.font.color.rgb = COLORS['text_medium']
    p.alignment = PP_ALIGN.CENTER

    # Feature cards
    card_width = Inches(2.6)
    card_height = Inches(2.2)
    start_x = Inches(0.7)
    gap = Inches(0.35)
    y = Inches(3.3)

    features = [
        ("24/7 Availability", "Always there when they need someone to talk to", COLORS['primary']),
        ("Personalized", "Conversations tailored to their interests and history", COLORS['secondary']),
        ("Family Updates", "Weekly summaries so you stay connected", COLORS['accent']),
    ]

    for i, (title, desc, accent) in enumerate(features):
        x = start_x + i * (card_width + gap)

        # Card background
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_width, card_height)
        card.fill.solid()
        card.fill.fore_color.rgb = COLORS['white']
        card.line.color.rgb = RgbColor(0xE5, 0xE7, 0xEB)
        card.line.width = Pt(1)

        # Accent bar at top
        accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, card_width, Inches(0.08))
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = accent
        accent_bar.line.fill.background()

        # Card title
        title_box = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.2), card_width - Inches(0.3), Inches(0.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "Merriweather"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLORS['text_dark']

        # Card description
        desc_box = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.7), card_width - Inches(0.3), Inches(1.2))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.name = "Inter"
        p.font.size = Pt(12)
        p.font.color.rgb = COLORS['text_medium']

    # CTA Button
    cta = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.5), Inches(5.8), Inches(3), Inches(0.7))
    cta.fill.solid()
    cta.fill.fore_color.rgb = COLORS['secondary']
    cta.line.fill.background()
    tf = cta.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = "Join the Beta"
    p.font.name = "Inter"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    tf.anchor = MSO_ANCHOR.MIDDLE

    return slide


def main():
    """Create the brand deck"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    script_dir = os.path.dirname(os.path.abspath(__file__))
    logo_path = os.path.join(script_dir, "logo.png")

    # Create slides
    add_title_slide(prs, "Lonesome No More", "Brand Guidelines for Gamma Import")
    add_logo_slide(prs, logo_path)
    add_color_palette_slide(prs)
    add_typography_slide(prs)
    add_button_styles_slide(prs)
    add_sample_content_slide(prs)
    add_title_slide(prs, "Thank You", "A family-owned labor of love • Built in Buffalo, NY")

    # Save
    output_path = os.path.join(script_dir, "lonesome-no-more-brand.pptx")
    prs.save(output_path)
    print(f"Brand deck created: {output_path}")
    return output_path


if __name__ == "__main__":
    main()
